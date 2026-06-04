#!/usr/bin/env python3
"""Create a stage-results PowerPoint deck for the CTV sparse-prompt project."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
FIG = ROOT / "manuscript_vsi_biomedical_data" / "figures"
OUT_DIR = ROOT / "reports" / "ppt"
PPT_ASSETS = OUT_DIR / "assets"
OUT_PATH = OUT_DIR / "CTV_sparse_prompt_stage_results_20260602.pptx"


COLORS = {
    "ink": RGBColor(26, 32, 44),
    "muted": RGBColor(89, 99, 116),
    "blue": RGBColor(32, 92, 160),
    "teal": RGBColor(12, 125, 116),
    "green": RGBColor(41, 130, 80),
    "red": RGBColor(174, 64, 64),
    "line": RGBColor(214, 219, 226),
    "pale": RGBColor(244, 247, 251),
    "white": RGBColor(255, 255, 255),
}


def set_text(run, size=22, bold=False, color="ink"):
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = COLORS[color]


def add_title(slide, title, subtitle=None):
    box = slide.shapes.add_textbox(Inches(0.55), Inches(0.3), Inches(12.2), Inches(0.65))
    p = box.text_frame.paragraphs[0]
    p.text = title
    p.alignment = PP_ALIGN.LEFT
    set_text(p.runs[0], 27, True, "ink")
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.58), Inches(0.92), Inches(11.8), Inches(0.35))
        p = sub.text_frame.paragraphs[0]
        p.text = subtitle
        set_text(p.runs[0], 12.5, False, "muted")
    line = slide.shapes.add_shape(1, Inches(0.55), Inches(1.22), Inches(12.2), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS["line"]
    line.line.color.rgb = COLORS["line"]


def add_bullets(slide, items, x, y, w, h, size=18, color="ink"):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.space_after = Pt(7)
        set_text(p.runs[0], size, False, color)
    return box


def add_card(slide, title, body, x, y, w, h, accent="blue"):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS["pale"]
    shape.line.color.rgb = COLORS["line"]
    shape.line.width = Pt(1)
    bar = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(0.09), Inches(h))
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLORS[accent]
    bar.line.color.rgb = COLORS[accent]
    box = slide.shapes.add_textbox(Inches(x + 0.18), Inches(y + 0.14), Inches(w - 0.32), Inches(h - 0.22))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    set_text(p.runs[0], 15, True, accent)
    for line in body:
        p = tf.add_paragraph()
        p.text = line
        p.space_before = Pt(2)
        set_text(p.runs[0], 13, False, "ink")


def add_flow_box(slide, idx, title, body, x, y, w, h, accent="blue"):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS["white"]
    shape.line.color.rgb = COLORS[accent]
    shape.line.width = Pt(1.2)

    strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.28))
    strip.fill.solid()
    strip.fill.fore_color.rgb = COLORS[accent]
    strip.line.color.rgb = COLORS[accent]
    text = strip.text_frame
    text.margin_left = Inches(0.06)
    text.margin_right = Inches(0.04)
    p = text.paragraphs[0]
    p.text = f"{idx}"
    p.alignment = PP_ALIGN.LEFT
    set_text(p.runs[0], 11, True, "white")

    box = slide.shapes.add_textbox(Inches(x + 0.12), Inches(y + 0.38), Inches(w - 0.24), Inches(h - 0.45))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.text = title
    set_text(p.runs[0], 13, True, accent)
    for line in body:
        p = tf.add_paragraph()
        p.text = line
        p.space_before = Pt(2)
        set_text(p.runs[0], 10.3, False, "ink")


def add_small_arrow(slide, x, y, w=0.34, h=0.22):
    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(w), Inches(h))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = COLORS["line"]
    arrow.line.color.rgb = COLORS["line"]


def add_method_schematic(slide):
    lane_w = 12.0
    lane = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.68), Inches(1.42), Inches(lane_w), Inches(2.78))
    lane.fill.solid()
    lane.fill.fore_color.rgb = COLORS["pale"]
    lane.line.color.rgb = COLORS["line"]
    lane.line.width = Pt(1)

    label = slide.shapes.add_textbox(Inches(0.9), Inches(1.48), Inches(2.6), Inches(0.25))
    p = label.text_frame.paragraphs[0]
    p.text = "Whole-volume preprocessing"
    set_text(p.runs[0], 12.5, True, "ink")

    add_flow_box(slide, 1, "Global CT", ["full thoracic volume", "single patient case"], 0.9, 1.82, 1.78, 1.0, "blue")
    add_flow_box(slide, 2, "OAR segmentation", ["lung / heart / cord", "esophagus / body"], 3.02, 1.82, 1.9, 1.0, "teal")
    add_flow_box(slide, 3, "Organ ROI", ["anatomy-aware search", "exclude impossible space"], 5.27, 1.82, 2.0, 1.0, "green")

    add_flow_box(slide, 4, "Sparse CTV prompt", ["K=7 clinician slices", "2D target masks"], 3.02, 3.0, 1.9, 1.0, "blue")
    add_flow_box(slide, 5, "SDF candidates", ["z propagation", "multi-variant support"], 5.27, 3.0, 2.0, 1.0, "teal")

    add_flow_box(slide, 6, "Core-envelope", ["C: high precision core", "E: high recall envelope"], 8.05, 2.22, 2.05, 1.3, "green")
    add_flow_box(slide, 7, "Final CTV", ["support-intersection rule", "3D pseudo label"], 10.72, 2.22, 1.85, 1.3, "red")

    add_small_arrow(slide, 2.73, 2.22, 0.26, 0.20)
    add_small_arrow(slide, 4.98, 2.22, 0.26, 0.20)
    add_small_arrow(slide, 4.98, 3.40, 0.26, 0.20)
    add_small_arrow(slide, 7.36, 2.50, 0.45, 0.24)
    add_small_arrow(slide, 7.36, 3.18, 0.45, 0.24)
    add_small_arrow(slide, 10.18, 2.86, 0.45, 0.24)

    bridge = slide.shapes.add_textbox(Inches(7.28), Inches(2.82), Inches(0.62), Inches(0.35))
    p = bridge.text_frame.paragraphs[0]
    p.text = "merge"
    p.alignment = PP_ALIGN.CENTER
    set_text(p.runs[0], 9.5, False, "muted")

    canvas = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.68), Inches(4.55), Inches(12.0), Inches(1.78))
    canvas.fill.solid()
    canvas.fill.fore_color.rgb = COLORS["white"]
    canvas.line.color.rgb = COLORS["line"]
    canvas.line.width = Pt(1)

    header = slide.shapes.add_textbox(Inches(0.95), Inches(4.72), Inches(4.0), Inches(0.28))
    p = header.text_frame.paragraphs[0]
    p.text = "CTV completion is restricted inside the organ ROI"
    set_text(p.runs[0], 13, True, "ink")

    roi = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.95), Inches(5.12), Inches(4.95), Inches(0.85))
    roi.fill.solid()
    roi.fill.fore_color.rgb = RGBColor(235, 242, 252)
    roi.line.color.rgb = COLORS["blue"]
    roi.line.width = Pt(1.2)
    p = roi.text_frame.paragraphs[0]
    p.text = "Organ ROI: anatomy-aware candidate space"
    p.alignment = PP_ALIGN.CENTER
    set_text(p.runs[0], 10.5, True, "blue")

    envelope = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.45), Inches(5.38), Inches(3.65), Inches(0.38))
    envelope.fill.solid()
    envelope.fill.fore_color.rgb = RGBColor(225, 243, 234)
    envelope.line.color.rgb = COLORS["green"]
    envelope.line.width = Pt(1.2)
    p = envelope.text_frame.paragraphs[0]
    p.text = "Envelope E"
    p.alignment = PP_ALIGN.CENTER
    set_text(p.runs[0], 9.5, True, "green")

    core = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2.12), Inches(5.48), Inches(1.35), Inches(0.18))
    core.fill.solid()
    core.fill.fore_color.rgb = COLORS["green"]
    core.line.color.rgb = COLORS["green"]

    uncertain = slide.shapes.add_textbox(Inches(5.12), Inches(5.26), Inches(1.0), Inches(0.55))
    p = uncertain.text_frame.paragraphs[0]
    p.text = "U = E - C"
    p.alignment = PP_ALIGN.CENTER
    set_text(p.runs[0], 10, False, "muted")

    formula = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.45), Inches(5.05), Inches(2.35), Inches(0.72))
    formula.fill.solid()
    formula.fill.fore_color.rgb = COLORS["pale"]
    formula.line.color.rgb = COLORS["line"]
    p = formula.text_frame.paragraphs[0]
    p.text = "C ⊂ Y_final ⊂ E"
    p.alignment = PP_ALIGN.CENTER
    set_text(p.runs[0], 15, True, "teal")

    note = slide.shapes.add_textbox(Inches(9.15), Inches(4.88), Inches(3.15), Inches(0.9))
    tf = note.text_frame
    tf.word_wrap = True
    for i, line in enumerate(["Train: calibrate support rule", "Inference: CT + OAR + sparse prompt", "Full GT: evaluation only"]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.space_after = Pt(3)
        set_text(p.runs[0], 10.5, False, "ink")

    add_bullets(
        slide,
        [
            "The method converts whole-volume CT into an anatomy-constrained CTV completion problem.",
            "OARs define the legal ROI; sparse CTV prompts define the target-specific SDF support.",
        ],
        0.95,
        6.48,
        11.6,
        0.55,
        12.5,
        "muted",
    )


def add_picture(slide, rel_path, x, y, w=None, h=None):
    path = FIG / rel_path
    if not path.exists():
        return None
    kwargs = {}
    if w is not None:
        kwargs["width"] = Inches(w)
    if h is not None:
        kwargs["height"] = Inches(h)
    return slide.shapes.add_picture(str(path), Inches(x), Inches(y), **kwargs)


def add_asset_picture(slide, filename, x, y, w=None, h=None):
    path = PPT_ASSETS / filename
    if not path.exists():
        return None
    kwargs = {}
    if w is not None:
        kwargs["width"] = Inches(w)
    if h is not None:
        kwargs["height"] = Inches(h)
    return slide.shapes.add_picture(str(path), Inches(x), Inches(y), **kwargs)


def add_formula_box(slide, title, lines, x, y, w, h, accent="blue", body_size=13):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS["white"]
    shape.line.color.rgb = COLORS["line"]
    shape.line.width = Pt(1)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.08), Inches(h))
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLORS[accent]
    bar.line.color.rgb = COLORS[accent]
    box = slide.shapes.add_textbox(Inches(x + 0.18), Inches(y + 0.12), Inches(w - 0.35), Inches(h - 0.18))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    set_text(p.runs[0], 13.5, True, accent)
    for line in lines:
        p = tf.add_paragraph()
        p.text = line
        p.space_before = Pt(2)
        set_text(p.runs[0], body_size, False, "ink")


def add_method_overview_visual(slide):
    add_asset_picture(slide, "method_overview_global_roi_ctv.png", 0.45, 1.35, w=12.35)
    add_formula_box(
        slide,
        "Pipeline",
        [
            "Global CT → OAR segmentation → organ ROI",
            "Sparse CTV slices → SDF candidates → CTV pseudo label",
        ],
        0.85,
        5.08,
        5.9,
        1.15,
        "teal",
        12.2,
    )
    add_formula_box(
        slide,
        "Inference / evaluation boundary",
        [
            "Inference uses CT + OAR + sparse CTV prompts only",
            "Complete CTV GT is hidden until final metric computation",
        ],
        7.0,
        5.08,
        5.5,
        1.15,
        "blue",
        12.2,
    )
    add_bullets(
        slide,
        ["This slide shows the full data-preprocessing path from whole CT to anatomy-constrained CTV completion."],
        0.9,
        6.52,
        11.6,
        0.35,
        12.2,
        "muted",
    )


def add_sdf_propagation_slide(slide):
    add_asset_picture(slide, "method_sdf_propagation_detail.png", 0.55, 1.42, w=8.1)
    add_asset_picture(slide, "method_prompt_z_profile.png", 9.0, 1.55, w=3.5)
    add_formula_box(
        slide,
        "SDF interpolation",
        [
            "D_k(x,y) = signed_distance(M_k)",
            "D(z,x,y) = Interp_z({D_k})",
            "Y_sdf = 1[D(z,x,y) ≤ 0]",
        ],
        0.75,
        4.72,
        5.9,
        1.45,
        "teal",
        12.3,
    )
    add_formula_box(
        slide,
        "Before → after",
        [
            "Before: only K sparse 2D CTV masks are known",
            "After: a dense 3D CTV candidate is generated by SDF propagation",
        ],
        6.95,
        4.72,
        5.55,
        1.45,
        "blue",
        12.3,
    )


def add_core_envelope_slide(slide):
    add_asset_picture(slide, "method_core_envelope_detail.png", 0.45, 1.38, w=12.35)
    add_formula_box(
        slide,
        "Support consensus",
        [
            "S(v) = (1/M) Σ_m 1[v ∈ Y_m]",
            "C = {v | S(v) ≥ τ_core}",
            "E = {v | S(v) ≥ τ_env},  τ_env ≤ τ_core",
        ],
        0.75,
        4.78,
        6.25,
        1.55,
        "green",
        12.0,
    )
    add_formula_box(
        slide,
        "Core-envelope prior",
        [
            "C: high-precision foreground that should be preserved",
            "E: high-recall candidate envelope that bounds the search space",
            "U = E - C: only uncertain region needs refinement",
        ],
        7.25,
        4.78,
        5.25,
        1.55,
        "teal",
        11.8,
    )


def add_support_rule_slide(slide):
    add_asset_picture(slide, "method_support_rule_detail.png", 0.45, 1.38, w=12.35)
    add_formula_box(
        slide,
        "Train-calibrated selection",
        [
            "r = |C_base| / |Y_base|",
            "Y_final = support_100, if r ≥ θ",
            "Y_final = linear_core_intersection, otherwise",
            "θ = 0.9909, selected on the training set",
        ],
        0.75,
        4.70,
        6.15,
        1.75,
        "red",
        11.5,
    )
    add_formula_box(
        slide,
        "Output constraint",
        [
            "C ⊂ Y_final ⊂ E",
            "The method changes only the candidate CTV support, not the whole CT volume",
            "Complete GT is used only after inference for Dice / HD95 / ASD",
        ],
        7.15,
        4.70,
        5.35,
        1.75,
        "blue",
        11.5,
    )


def add_visual_only_slide(slide, filename, caption, y=1.45, w=12.25):
    add_asset_picture(slide, filename, 0.55, y, w=w)
    add_bullets(
        slide,
        [caption],
        0.85,
        6.58,
        11.8,
        0.35,
        12.2,
        "muted",
    )


def add_table(slide, rows, cols, x, y, w, h, headers, data, font_size=12):
    table = slide.shapes.add_table(rows, cols, Inches(x), Inches(y), Inches(w), Inches(h)).table
    for i, width in enumerate([w / cols] * cols):
        table.columns[i].width = Inches(width)
    for c, text in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = text
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLORS["blue"]
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        set_text(p.runs[0], font_size, True, "white")
    for r, row in enumerate(data, start=1):
        for c, text in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(text)
            cell.fill.solid()
            cell.fill.fore_color.rgb = COLORS["white"] if r % 2 else COLORS["pale"]
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if c > 1 else PP_ALIGN.LEFT
            set_text(p.runs[0], font_size, False, "ink")
    return table


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def build_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 1
    slide = blank_slide(prs)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLORS["white"]
    box = slide.shapes.add_textbox(Inches(0.72), Inches(1.0), Inches(11.9), Inches(1.15))
    p = box.text_frame.paragraphs[0]
    p.text = "稀疏提示驱动的 CTV 伪标签生成"
    set_text(p.runs[0], 34, True, "ink")
    box = slide.shapes.add_textbox(Inches(0.75), Inches(2.08), Inches(11.0), Inches(0.55))
    p = box.text_frame.paragraphs[0]
    p.text = "阶段性成果汇报：data preprocessing / sparse annotation completion"
    set_text(p.runs[0], 20, False, "blue")
    add_card(slide, "核心结论", ["K=7 sparse CTV prompts + SDF support-intersection", "测试集 Dice 0.928，Unseen Dice 0.897", "任务定位：生成高质量 3D CTV pseudo label"], 0.8, 3.15, 5.8, 1.55, "teal")
    add_card(slide, "当前状态", ["主实验与论文包已完成", "无 GTV 下游 nnU-Net 数据集已准备", "GPU driver 当前阻塞，训练待服务器恢复"], 6.85, 3.15, 5.6, 1.55, "blue")
    add_bullets(slide, ["Dataset: 34 train / 31 test", "Target: CTV completion from sparse full-slice prompts", "Current date: 2026-06-02"], 0.8, 5.45, 11.5, 0.8, 14, "muted")

    # 2
    slide = blank_slide(prs)
    add_title(slide, "问题定位", "CTV 不是纯 CT-visible structure，直接全自动分割不是本文主线")
    add_card(slide, "不定义为", ["CT-only automatic CTV segmentation", "全图自由预测 CTV mask", "依赖网络猜测医生临床意图"], 0.7, 1.6, 3.9, 2.0, "red")
    add_card(slide, "定义为", ["少量 2D CTV 标注作为 sparse prompt", "OAR/解剖约束限制搜索空间", "SDF core-envelope 补全 3D pseudo label"], 4.75, 1.6, 4.1, 2.0, "teal")
    add_card(slide, "论文贡献", ["将稀疏标注预处理为完整 3D CTV 标签", "降低完整 3D 勾画成本", "为后续训练提供更可靠监督"], 9.0, 1.6, 3.6, 2.0, "blue")
    add_picture(slide, "vsi_dataset_ctv_concept.png", 1.05, 4.0, w=11.2)

    # 3
    slide = blank_slide(prs)
    add_title(slide, "数据与划分", "病例级划分，测试时完整 GT 仅用于最终评价")
    add_table(
        slide,
        5,
        4,
        0.75,
        1.55,
        11.9,
        1.8,
        ["Split / Dataset", "Cases", "Role", "Notes"],
        [
            ["Train", "34", "校准规则 / 训练网络", "每例一个 CTV case"],
            ["Test", "31", "最终评价", "GT hidden during inference"],
            ["Pseudo train", "Dataset016-018", "下游监督", "Linear / SDF core / Ours"],
            ["Full GT", "Dataset015", "Upper bound", "真实标签训练 nnU-Net"],
        ],
        11,
    )
    add_bullets(
        slide,
        [
            "当前主方法推理不需要 GTV；GTV 可作为未来 disease-origin prior。",
            "下游任务使用 CT-only 输入，唯一变量是训练标签来源。",
            "伪标签训练集已完成 raw + nnU-Net preprocess。"
        ],
        0.95,
        4.0,
        11.2,
        1.3,
        17,
    )
    add_table(
        slide,
        4,
        4,
        1.35,
        5.5,
        10.2,
        1.05,
        ["Dataset", "Label source", "Train n", "Pseudo Dice to GT"],
        [
            ["016", "Linear", "34", "0.923 ± 0.047"],
            ["017", "SDF core", "34", "0.928 ± 0.042"],
            ["018", "Ours", "34", "0.936 ± 0.045"],
        ],
        11,
    )

    # 4
    slide = blank_slide(prs)
    add_title(slide, "方法总流程", "Global CT → OAR-based ROI → sparse-prompt CTV completion")
    add_method_overview_visual(slide)

    # 5
    slide = blank_slide(prs)
    add_title(slide, "SDF Propagation", "从少量 2D CTV prompt 生成连续 3D CTV 候选")
    add_sdf_propagation_slide(slide)

    # 6
    slide = blank_slide(prs)
    add_title(slide, "Core-Envelope Prior", "用多候选支持度构建高精度 core 与高召回 envelope")
    add_core_envelope_slide(slide)

    # 7
    slide = blank_slide(prs)
    add_title(slide, "Support Rule", "训练集校准的确定性规则完成最终 CTV 选择")
    add_support_rule_slide(slide)

    # 8
    slide = blank_slide(prs)
    add_title(slide, "全流程可视化", "从全局 CT 到器官 ROI，再到 CTV 伪标签输出")
    add_visual_only_slide(
        slide,
        "method_overview_global_roi_ctv.png",
        "展示全局 3D CT、OAR 局部 ROI、稀疏 CTV prompt 和最终 CTV pseudo label 的完整预处理链路。",
        1.38,
        12.25,
    )

    # 9
    slide = blank_slide(prs)
    add_title(slide, "SDF Propagation 可视化", "Sparse 2D prompt → SDF field → dense 3D candidate")
    add_visual_only_slide(
        slide,
        "method_sdf_propagation_detail.png",
        "该页单独展示 SDF propagation 的前后变化：稀疏 2D 标注被转换为连续 SDF 场，并生成 dense CTV candidate。",
        1.62,
        11.9,
    )

    # 10
    slide = blank_slide(prs)
    add_title(slide, "Core-Envelope 可视化", "从 SDF support 构建 high-precision core 与 high-recall envelope")
    add_visual_only_slide(
        slide,
        "method_core_envelope_detail.png",
        "该页单独展示 core-envelope prior：core 保留高置信前景，envelope 限定候选搜索空间，U=E-C 是需要处理的不确定区域。",
        1.55,
        12.25,
    )

    # 11
    slide = blank_slide(prs)
    add_title(slide, "Support Rule 可视化", "Linear / SDF candidate → rule-selected final CTV")
    add_visual_only_slide(
        slide,
        "method_support_rule_detail.png",
        "该页单独展示 support rule 的处理前后变化，并强调最终输出仍受 C ⊂ Y_final ⊂ E 约束。",
        1.55,
        12.25,
    )

    # 12
    slide = blank_slide(prs)
    add_title(slide, "Baseline 结果说明", "CT-only 与普通 prompt baseline 说明 CTV 任务难度")
    add_table(
        slide,
        6,
        5,
        0.55,
        1.45,
        12.2,
        2.6,
        ["Group", "Method", "Dice", "HD95", "ASD"],
        [
            ["No prompt", "nnU-Net", "0.400 ± 0.345", "56.94", "33.72"],
            ["No prompt", "DiffUNet", "0.316 ± 0.301", "71.89", "41.44"],
            ["Auto prompt", "SAM-Med3D CT-derived", "0.157 ± 0.187", "85.70", "49.64"],
            ["Sparse prompt", "SAM-Med3D K=7", "0.422 ± 0.159", "21.32", "5.54"],
            ["Sparse prompt", "Linear interpolation K=7", "0.912 ± 0.048", "3.90", "1.22"],
        ],
        10.5,
    )
    add_picture(slide, "baseline_ctv_overlay.png", 0.95, 4.45, w=5.65)
    add_picture(slide, "sammed3d_sparse_prompt_k7_ctv_overlay.png", 6.75, 4.45, w=5.65)

    # 6
    slide = blank_slide(prs)
    add_title(slide, "主结果：K=7 CTV 补全", "Ours 在同样 sparse full-slice prompt 条件下优于传统插值和 SDF core")
    add_table(
        slide,
        6,
        6,
        0.45,
        1.35,
        12.45,
        2.35,
        ["Method", "Dice", "Unseen Dice", "HD95", "ASD", "Vol diff"],
        [
            ["Linear interpolation", "0.912 ± 0.048", "0.876 ± 0.057", "3.90", "1.22", "1.4%"],
            ["SDF core", "0.920 ± 0.045", "0.879 ± 0.061", "4.21", "1.13", "-3.7%"],
            ["Strict support", "0.923 ± 0.046", "0.885 ± 0.059", "4.00", "1.07", "-4.3%"],
            ["Linear-core intersection", "0.920 ± 0.054", "0.886 ± 0.068", "3.68", "1.10", "-9.5%"],
            ["Ours support-intersection", "0.928 ± 0.050", "0.897 ± 0.062", "3.43", "0.97", "-6.6%"],
        ],
        10.2,
    )
    add_picture(slide, "vsi_main_results_dice.png", 1.0, 4.05, w=11.0)

    # 7
    slide = blank_slide(prs)
    add_title(slide, "配对统计", "Support-intersection rule 对 Linear / SDF core 均有显著提升")
    add_picture(slide, "vsi_paired_dice_delta.png", 0.75, 1.35, w=7.3)
    add_card(slide, "Scan-level", ["vs Linear: ΔDice +0.0160", "Improved/Worse: 26/5", "p = 5.32e-06"], 8.35, 1.55, 4.0, 1.65, "teal")
    add_card(slide, "Patient-level", ["vs Linear: ΔDice +0.0169", "Improved/Worse: 19/2", "p = 1.61e-04"], 8.35, 3.45, 4.0, 1.65, "blue")
    add_card(slide, "vs SDF core", ["Scan: ΔDice +0.0081, p=2.88e-06", "Patient: ΔDice +0.0089, p=8.86e-05"], 8.35, 5.35, 4.0, 1.2, "green")

    # 8
    slide = blank_slide(prs)
    add_title(slide, "结构化消融", "最终结果来自 SDF support consensus 与 linear-core intersection 的组合")
    add_picture(slide, "vsi_structured_ablation_refinement.png", 0.65, 1.35, w=12.0)
    add_bullets(slide, ["Doctor-prior graph refinement 未形成正向主结果，保留为负结果/未来方向。", "当前 Ours 是 deterministic preprocessing rule，不是额外网络。"], 0.85, 6.55, 11.8, 0.5, 13.2, "muted")

    # 9
    slide = blank_slide(prs)
    add_title(slide, "临床鲁棒性", "绝大多数测试病例通过固定 review gates")
    add_picture(slide, "vsi_clinical_robustness_gates.png", 0.65, 1.35, w=7.4)
    add_table(
        slide,
        6,
        3,
        8.2,
        1.6,
        4.2,
        3.2,
        ["Gate", "Pass", "Meaning"],
        [
            ["Dice ≥ 0.85", "29/31", "overall overlap"],
            ["Unseen Dice ≥ 0.80", "28/31", "non-prompt slices"],
            ["HD95 ≤ 6 mm", "29/31", "surface outliers"],
            ["ASD ≤ 2 mm", "31/31", "mean surface error"],
            ["|ΔV|≤15%", "28/31", "volume bias"],
        ],
        10.5,
    )

    # 10
    slide = blank_slide(prs)
    add_title(slide, "定性对照", "同一病例中，Ours 更接近完整 CTV 范围")
    add_picture(slide, "our_sdf_k7_ctv_main_comparison.png", 0.7, 1.35, w=12.0)
    add_bullets(slide, ["图中重点：prompt slices 之外的补全质量，以及 envelope 内过扩/欠扩控制。"], 0.9, 6.65, 11.5, 0.4, 13, "muted")

    # 11
    slide = blank_slide(prs)
    add_title(slide, "OAR 的角色", "OAR 是上游辅助条件，不是本文核心输出")
    add_picture(slide, "vsi_oar_segmentation_dice.png", 0.75, 1.35, w=7.0)
    add_card(slide, "OAR 在流程中的用途", ["构建 anatomy-aware search space", "限制不合理扩散", "作为距离/解剖先验"], 8.0, 1.65, 4.4, 1.9, "blue")
    add_card(slide, "论文表述", ["OAR segmentation is auxiliary", "CTV sparse-prompt preprocessing is the core", "不把 OAR baseline comparison 作为主贡献"], 8.0, 3.9, 4.4, 1.9, "teal")

    # 12
    slide = blank_slide(prs)
    add_title(slide, "无 GTV 下游任务设计", "目标是验证伪标签作为训练监督的 utility，而不是让网络超过伪标签本身")
    add_table(
        slide,
        5,
        5,
        0.6,
        1.45,
        12.0,
        2.35,
        ["Experiment", "Input", "Training label", "Dataset", "Role"],
        [
            ["A", "CT", "Linear pseudo", "016", "weak pseudo baseline"],
            ["B", "CT", "SDF core pseudo", "017", "structured pseudo baseline"],
            ["C", "CT", "Ours pseudo", "018", "proposed supervision"],
            ["D", "CT", "Full GT", "015", "upper bound / reference"],
        ],
        10.5,
    )
    add_card(slide, "Expected claim", ["Ours pseudo labels provide better downstream supervision than competing pseudo-label preprocessing strategies."], 0.9, 4.35, 5.7, 1.25, "green")
    add_card(slide, "Do not claim", ["Pseudo-label-trained CT-only nnU-Net must exceed the pseudo label itself or full-GT training."], 6.95, 4.35, 5.45, 1.25, "red")
    add_bullets(slide, ["如果 CT-only downstream 仍低，结论是 CTV 需要 case-specific sparse clinician prompts。"], 0.9, 6.25, 11.5, 0.5, 14, "muted")

    # 13
    slide = blank_slide(prs)
    add_title(slide, "当前工程状态", "下游数据准备完成；GPU driver 阻塞导致训练暂缓")
    add_card(slide, "已完成", ["Dataset016/017/018 raw datasets", "nnU-Net 3d_fullres preprocess", "阶段性论文图表与可视化"], 0.85, 1.6, 4.0, 2.15, "green")
    add_card(slide, "阻塞", ["GPU3: Unknown Error", "新 CUDA 初始化进入 D/Dl 状态", "训练日志 0 字节，未进入显存"], 4.95, 1.6, 4.0, 2.15, "red")
    add_card(slide, "恢复后命令", ["bash scripts/run_downstream_pseudo_ctv_nnunet_server05.sh train", "默认跳过 GPU3，先跑 016/017/018 fold0"], 9.05, 1.6, 3.55, 2.15, "blue")
    add_bullets(
        slide,
        [
            "脚本已增加保护：发现 D-state CUDA 残留时拒绝继续提交。",
            "preflight 已取消 torch.cuda.is_available()，避免预检本身卡死。",
            "后续训练完成后，再统一 predict + evaluate 31 test cases。"
        ],
        0.95,
        4.5,
        11.7,
        1.5,
        16,
    )

    # 14
    slide = blank_slide(prs)
    add_title(slide, "阶段结论与下一步", "当前工作足以支撑 data preprocessing 主线，下一步补 utility validation")
    add_card(slide, "论文主结论", ["稀疏 2D CTV prompts 可以被预处理为高质量 3D CTV pseudo labels", "Ours support-intersection rule 显著优于 linear 与 SDF core"], 0.8, 1.55, 5.8, 1.75, "teal")
    add_card(slide, "需要补充", ["无 GTV downstream utility validation", "训练完成后生成 test metrics 和三线表", "可选：GTV 作为 future extension / ablation"], 6.9, 1.55, 5.6, 1.75, "blue")
    add_bullets(
        slide,
        [
            "短期：重启/修复 GPU driver 后，启动 Dataset016/017/018 fold0 训练。",
            "中期：完成同一 31 test cases 上的 predict/evaluate。",
            "写作：下游任务定位为 utility validation，不作为主 refinement network。"
        ],
        1.0,
        4.2,
        11.3,
        1.55,
        18,
    )

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(OUT_PATH)
    return OUT_PATH


if __name__ == "__main__":
    path = build_deck()
    print(path)
